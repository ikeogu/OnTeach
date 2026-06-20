"""Extract plain text from PDF, DOCX, PPTX, TXT files — local paths or remote URLs."""
import tempfile
from pathlib import Path
from urllib.parse import urlparse


def extract_text(file_path_or_url: str) -> str:
    if file_path_or_url.startswith(("http://", "https://")):
        return _extract_from_url(file_path_or_url)

    path = Path(file_path_or_url)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {file_path_or_url}")
    return _extract_by_path(path)


def _extract_from_url(url: str) -> str:
    import httpx

    suffix = Path(urlparse(url).path).suffix.lower() or ".tmp"

    with httpx.Client(timeout=60, follow_redirects=True) as client:
        resp = client.get(url)
        resp.raise_for_status()
        content_type = resp.headers.get("content-type", "")

    # Web page — extract visible text from HTML
    if "text/html" in content_type or suffix in {".html", ".htm", ".tmp"}:
        return _extract_html(resp.content)

    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(resp.content)
        tmp_path = Path(tmp.name)

    try:
        return _extract_by_path(tmp_path)
    finally:
        tmp_path.unlink(missing_ok=True)


def _extract_html(content: bytes) -> str:
    import re

    text = content.decode("utf-8", errors="ignore")

    # Strip <script>, <style>, <nav>, <footer>, <header> blocks
    for tag in ("script", "style", "nav", "footer", "header"):
        text = re.sub(rf"<{tag}[^>]*>.*?</{tag}>", "", text, flags=re.DOTALL | re.IGNORECASE)

    # Remove all remaining HTML tags
    text = re.sub(r"<[^>]+>", " ", text)

    # Decode common HTML entities
    text = (
        text.replace("&amp;", "&")
        .replace("&lt;", "<")
        .replace("&gt;", ">")
        .replace("&nbsp;", " ")
        .replace("&#39;", "'")
        .replace("&quot;", '"')
    )

    # Collapse whitespace
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _extract_by_path(path: Path) -> str:
    suffix = path.suffix.lower()

    if suffix in {".txt", ".md"}:
        return path.read_text(errors="ignore")

    if suffix == ".pdf":
        return _extract_pdf(path)

    if suffix in {".docx", ".doc"}:
        return _extract_docx(path)

    if suffix in {".pptx", ".ppt"}:
        return _extract_pptx(path)

    return path.read_text(errors="ignore")


def _extract_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    parts: list[str] = []
    for page in reader.pages:
        text = page.extract_text() or ""
        if text.strip():
            parts.append(text)
    return "\n\n".join(parts)


def _extract_docx(path: Path) -> str:
    import docx

    doc = docx.Document(str(path))
    return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        slide_texts.append(text)
        if slide_texts:
            parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))
    return "\n\n".join(parts)
